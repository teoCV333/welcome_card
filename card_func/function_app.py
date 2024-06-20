import io
import logging
import os
import uuid
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont
from datetime import datetime

def sanitize_input(input_string):
    """
    Sanitiza la entrada para prevenir ataques de inyección.
    """
    return input_string.replace("..", "").replace("\\", "")

def validate_file_path(file_path):
    """
    Valida la ruta del archivo para asegurar que es una ruta permitida.
    """
    # Aquí puedes añadir lógica para verificar que la ruta esté en un directorio permitido
    return os.path.abspath(file_path)

def main(file_path, text_to_replace, new_text):

    if not file_path or not text_to_replace or not new_text:
        return "Por favor especifica la ruta del archivo, el texto a reemplazar y el nuevo texto"

    # Sanitizar entradas
    file_path = sanitize_input(file_path)
    text_to_replace = sanitize_input(text_to_replace)
    new_text = sanitize_input(new_text)

    # Validar y obtener la ruta completa del archivo
    try:
        file_path = validate_file_path(file_path)
    except Exception as e:
        logging.error(f"Invalid file path: {e}")
        return f"Error en la ruta del archivo: {e}"

    # Leer el archivo PowerPoint desde el directorio local
    try:
        ppt = Presentation(file_path)
    except Exception as e:
        logging.error(f"Error al leer el archivo PowerPoint: {e}")
        return f"Error al leer el archivo PowerPoint: {e}"

    # Reemplazar texto en las diapositivas
    for slide in ppt.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if text_to_replace == run.text:
                            run.text = run.text.replace(text_to_replace, new_text)
                            run.font.bold = True  # Establecer negrilla
                            run.font.color.rgb = RGBColor(187, 207, 0)

    # Guardar cada diapositiva como imagen JPEG
    image_files = []
    timestamp = datetime.now().strftime("%Y%m%d%H%M%S")
    for i, slide in enumerate(ppt.slides):
        try:
            image = slide_to_image(slide, ppt.slide_width, ppt.slide_height)
            unique_filename = f"slide_{i}_{uuid.uuid4().hex}.jpg"
            image.save(unique_filename, 'JPEG')
            image_files.append(unique_filename)
            logging.info(f"Imagen guardada en: {unique_filename}")
        except Exception as e:
            logging.error(f"Error al guardar la imagen: {e}")
            return f"Error al guardar la imagen: {e}"

    logging.info("Proceso completado, las imágenes se han guardado en el directorio local")
    return "Proceso completado, las imágenes se han guardado en el directorio local"

def slide_to_image(slide, slide_width, slide_height):
    # Convertir las dimensiones de puntos a píxeles
    dpi = 96 
    width_px = int(slide_width * dpi / Pt(72))
    height_px = int(slide_height * dpi / Pt(72))

    # Crear una imagen en blanco con las dimensiones adecuadas
    img = Image.new('RGB', (width_px, height_px), color=(255, 255, 255))
    draw = ImageDraw.Draw(img)

    # Renderizar el contenido de la diapositiva en la imagen
    for shape in slide.shapes:
        if shape.shape_type == 13:  # 13 indica una imagen
            image_stream = shape.image.blob
            image = Image.open(io.BytesIO(image_stream))
            left = int(shape.left * dpi / Pt(72))
            top = int(shape.top * dpi / Pt(72))
            width = int(shape.width * dpi / Pt(72))
            height = int(shape.height * dpi / Pt(72))
            image = image.resize((width, height), Image.LANCZOS)
            img.paste(image, (left, top))
        elif shape.has_text_frame:
            left = int(shape.left * dpi / Pt(72))
            top = int(shape.top * dpi / Pt(72))
            width = int(shape.width * dpi / Pt(72))
            height = int(shape.height * dpi / Pt(72))
            text_frame_left = left + int(shape.text_frame.margin_left * dpi / Pt(72))
            text_frame_top = top + int(shape.text_frame.margin_top * dpi / Pt(72))

            for paragraph in shape.text_frame.paragraphs:
                text_height = 0  # Altura total del texto en el cuadro de texto
                line_spacing = 1.2  # Espacio entre líneas, ajustable según sea necesario

                # Calcular la altura total del texto en el cuadro de texto
                for run in paragraph.runs:
                    if run.text == new_text:
                        font = run.font
                        font_size = font.size.pt  # Tamaño de la fuente en puntos
                        font = ImageFont.truetype("arial.ttf", size=int(font_size))  # Cargar la fuente Arial con el tamaño correcto
                        bbox = draw.textbbox((0, 0), run.text, font=font)
                        text_width = bbox[2] - bbox[0]
                        text_height = bbox[3] - bbox[1]
                        text_height = int(text_height * line_spacing)
                        break

                # Calcular la posición vertical inicial para centrar el texto
                text_y = text_frame_top + (height - text_height) / 2

                for run in paragraph.runs:
                    if run.text == new_text:
                        text = run.text
                        font = run.font
                        font.bold = True
                        font_size = 48  # Tamaño de la fuente en puntos
                        font = ImageFont.truetype("arialbd.ttf", size=int(font_size))  # Cargar la fuente Arial con el tamaño correcto
                        bbox = draw.textbbox((text_frame_left, text_y), text, font=font)
                        text_width = bbox[2] - bbox[0]
                        text_x = text_frame_left + (width - text_width) / 2
                        draw.text((text_x, text_y), text, fill=(187, 207, 0), font=font)
                        text_y += text_height  # Avanzar la posición vertical para la siguiente línea de texto

    return img

if __name__ == "__main__":
    # Variables para prueba
    file_path = "path"
    text_to_replace = "text to replace"
    new_text = "new text"

    result = main(file_path, text_to_replace, new_text)
    print(result)
